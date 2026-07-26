import os
import sys
import copy
import re
import tempfile
import win32com.client
from pptx import Presentation

def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        # Full text replacement for paragraph
        full_text = p.text
        if not full_text:
            continue
        
        modified_text = full_text
        for key, val in replacements.items():
            pattern = re.compile(rf"<<\s*{re.escape(key)}\s*>>", re.IGNORECASE)
            modified_text = pattern.sub(str(val), modified_text)
        
        if modified_text != full_text:
            # If paragraph runs exist, update run 0 and clear rest to preserve primary font styling
            if p.runs:
                p.runs[0].text = modified_text
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = modified_text

def test_native_pptx_generation():
    pptx_path = r"C:\Users\ROHITH P\Downloads\WORKSHOP.pptx"
    if not os.path.exists(pptx_path):
        print("WORKSHOP.pptx not found!")
        return

    output_dir = os.path.abspath("test_output_certs")
    os.makedirs(output_dir, exist_ok=True)

    test_row = {
        "Name": "Preethika sri K",
        "Roll Number": "25CSR220",
        "Roll Number ": "25CSR220",
        "Date": "25.07.2026",
        "issue_date": "25.07.2026",
        "event_date": "25.07.2026",
        "event_name": "Machine Learning Workshop"
    }

    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    for shape in slide.shapes:
        replace_text_in_shape(shape, test_row)

    temp_pptx = os.path.join(output_dir, "temp_preethika.pptx")
    prs.save(temp_pptx)
    print("Saved modified PPTX to:", temp_pptx)

    # Convert to PDF using PowerPoint COM Application
    pdf_out = os.path.join(output_dir, "Preethika_sri_K_Certificate.pdf")
    
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    # ppSaveAsPDF = 32
    # WithWindow = False (0)
    pres = ppt_app.Presentations.Open(os.path.abspath(temp_pptx), WithWindow=False)
    pres.SaveAs(os.path.abspath(pdf_out), 32)
    pres.Close()
    ppt_app.Quit()

    print("Successfully generated native PDF:", pdf_out)
    if os.path.exists(pdf_out):
        print("PDF file size (bytes):", os.path.getsize(pdf_out))

if __name__ == "__main__":
    test_native_pptx_generation()
