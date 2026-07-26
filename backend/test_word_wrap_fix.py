import os
from pptx import Presentation
from generator import generate_single_native_pdf, get_active_pptx_template_path, replace_tokens_in_pptx_slide

def replace_tokens_with_word_wrap(slide, replacements: dict):
    """Replaces tokens while setting word_wrap=True to prevent text smudging and character overlap."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Enable Word Wrap on text frames so dynamic text flows naturally without squishing
            shape.text_frame.word_wrap = True
    
    replace_tokens_in_pptx_slide(slide, replacements)

template_path = get_active_pptx_template_path()
output_pdf = os.path.abspath(os.path.join(os.path.dirname(__file__), "generated_pdfs", "Test_No_Smudge.pdf"))

replacements = {
    "Name": "Preethika sri K",
    "Roll Number": "25CSR220",
    "Roll Number ": "25CSR220",
    "Date": "25.07.2026",
    "Date ": "25.07.2026",
    "event_name": "DATASET TO DECISION Workshop"
}

prs = Presentation(template_path)
slide = prs.slides[0]
replace_tokens_with_word_wrap(slide, replacements)

temp_pptx = os.path.abspath(os.path.join(os.path.dirname(__file__), "generated_pdfs", "temp_nosmudge.pptx"))
prs.save(temp_pptx)

generate_single_native_pdf(temp_pptx, replacements, output_pdf)
print(f"Generated No-Smudge PDF at: {output_pdf}")
