import os
from pptx import Presentation

pptx_path = r"C:\Users\ROHITH P\Downloads\WORKSHOP.pptx"
if not os.path.exists(pptx_path):
    pptx_path = os.path.abspath(os.path.join(os.path.dirname(__file__), 'uploaded_templates', 'latest_template.pptx'))

print(f"Inspecting PPTX: {pptx_path}")
prs = Presentation(pptx_path)

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n--- SLIDE {slide_idx+1} ---")
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"\nShape {shape_idx+1} (TextFrame):")
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                print(f"  Paragraph {p_idx+1} (Full Text: '{p.text}'):")
                for r_idx, r in enumerate(p.runs):
                    font_name = r.font.name if r.font else 'Default'
                    font_size = r.font.size.pt if (r.font and r.font.size) else 'Default'
                    bold = r.font.bold if r.font else 'Default'
                    color = r.font.color.rgb if (r.font and r.font.color and hasattr(r.font.color, 'rgb')) else 'Default'
                    print(f"    Run {r_idx+1}: '{r.text}' | Font: {font_name}, Size: {font_size}, Bold: {bold}, Color: {color}")
