import os
import re
from pptx import Presentation
from pptx.dml.color import RGBColor

def replace_tokens_in_pptx_slide_enforce_blue_bold(slide, replacements: dict):
    """Replaces tokens while explicitly enforcing Bold + Blue (1A3987) for event_name / DATASET TO DECISION."""
    blue_color = RGBColor(0x1A, 0x39, 0x87)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        try:
            shape.text_frame.word_wrap = True
        except Exception:
            pass

        for p in shape.text_frame.paragraphs:
            full_text = p.text
            if not full_text:
                continue

            for r in p.runs:
                if not r.text:
                    continue

                # 1. If run contains DATASET TO DECISION or event_name, enforce Bold & Blue (1A3987)
                if 'DATASET TO DECISION' in r.text or '<<event_name>>' in r.text.lower() or '<<event name>>' in r.text.lower():
                    if r.font:
                        r.font.bold = True
                        r.font.color.rgb = blue_color

                # 2. Perform token replacement inside run
                if '<<' in r.text:
                    run_text = r.text
                    for key, val in replacements.items():
                        if not key:
                            continue
                        pattern = re.compile(rf"<<\s*{re.escape(str(key))}\s*>>", re.IGNORECASE)
                        run_text = pattern.sub(str(val), run_text)
                    r.text = run_text
                
                # Double-check post-replacement if event name is inside run
                event_val = str(replacements.get('event_name') or replacements.get('Event Name') or 'DATASET TO DECISION').strip()
                if event_val and event_val.lower() in r.text.lower() and len(event_val) > 2:
                    if r.font:
                        r.font.bold = True
                        r.font.color.rgb = blue_color

# Test execution
pptx_path = r"C:\Users\ROHITH P\Downloads\WORKSHOP.pptx"
prs = Presentation(pptx_path)
slide = prs.slides[0]

replacements = {
    "Name": "Preethika sri K",
    "Roll Number": "25CSR220",
    "Roll Number ": "25CSR220",
    "Date": "25.07.2026",
    "Date ": "25.07.2026",
    "event_name": "DATASET TO DECISION"
}

replace_tokens_in_pptx_slide_enforce_blue_bold(slide, replacements)

for shape_idx, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if 'DATASET' in p.text:
                print(f"Paragraph: '{p.text}'")
                for r in p.runs:
                    color = r.font.color.rgb if (r.font and r.font.color and hasattr(r.font.color, 'rgb')) else 'Default'
                    bold = r.font.bold if r.font else 'Default'
                    print(f"  Run '{r.text}' | Bold: {bold} | Color: {color}")
