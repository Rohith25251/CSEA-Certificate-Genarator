import sys
import os
from pptx import Presentation

pptx_path = r"C:\Users\ROHITH P\Downloads\WORKSHOP.pptx"
if not os.path.exists(pptx_path):
    print("File not found:", pptx_path)
    sys.exit(1)

prs = Presentation(pptx_path)
print("Slide width (pt):", prs.slide_width.pt, "height (pt):", prs.slide_height.pt)
slide = prs.slides[0]

for idx, shape in enumerate(slide.shapes):
    print(f"\n--- Shape {idx}: {shape.name} (type: {shape.shape_type}) ---")
    print(f"Position (pt): left={shape.left.pt:.1f}, top={shape.top.pt:.1f}, width={shape.width.pt:.1f}, height={shape.height.pt:.1f}")
    if shape.has_text_frame:
        for p_idx, p in enumerate(shape.text_frame.paragraphs):
            print(f"  P{p_idx} (align={p.alignment}): '{p.text}'")
            for r_idx, r in enumerate(p.runs):
                font_sz = r.font.size.pt if r.font.size else "N/A"
                print(f"    R{r_idx} (font={r.font.name}, sz={font_sz}, bold={r.font.bold}): '{r.text}'")
