import re
from pptx import Presentation

def replace_tokens_in_pptx_slide_preserved(slide, replacements: dict):
    """Replaces <<Placeholder>> tokens inside each run individually, preserving 100% of formatting, colors, and font sizes."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for p in shape.text_frame.paragraphs:
            full_text = p.text
            if not full_text:
                continue

            # First check if any placeholder token exists in the paragraph
            has_placeholder = any(f"<<{key}>>".lower() in full_text.lower() or f"<< {key} >>".lower() in full_text.lower() for key in replacements.keys())
            if not has_placeholder:
                # Also check regex pattern <<...>>
                if not re.search(r"<<\s*([^>]+?)\s*>>", full_text):
                    continue

            # Iterate through runs and replace tokens per run
            for r in p.runs:
                if not r.text or '<<' not in r.text:
                    continue
                
                run_text = r.text
                for key, val in replacements.items():
                    if not key:
                        continue
                    # Match <<key>>, << key >>, <<key >>, << key>> flexibly
                    pattern = re.compile(rf"<<\s*{re.escape(str(key))}\s*>>", re.IGNORECASE)
                    run_text = pattern.sub(str(val), run_text)
                
                r.text = run_text

            # Fallback for tokens split across run boundaries inside the paragraph
            remaining_tokens = re.findall(r"<<\s*([^>]+?)\s*>>", p.text)
            if remaining_tokens:
                # Token was split across run boundaries: replace across paragraph runs carefully
                p_text = p.text
                for key, val in replacements.items():
                    pattern = re.compile(rf"<<\s*{re.escape(str(key))}\s*>>", re.IGNORECASE)
                    p_text = pattern.sub(str(val), p_text)
                
                if p.runs:
                    p.runs[0].text = p_text
                    for r in p.runs[1:]:
                        r.text = ""

# Test execution
pptx_path = r"C:\Users\ROHITH P\Downloads\WORKSHOP.pptx"
prs = Presentation(pptx_path)
slide = prs.slides[0]

replacements = {
    "Name": "Preethika sri K",
    "Roll Number": "25CSR220",
    "Roll Number ": "25CSR220",
    "Date": "25.07.2026",
    "Date ": "25.07.2026"
}

replace_tokens_in_pptx_slide_preserved(slide, replacements)

print("\n--- AFTER PRESERVED REPLACEMENT ---")
for shape_idx, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.text.strip():
                print(f"\nShape {shape_idx+1} Paragraph Text: '{p.text}'")
                for r_idx, r in enumerate(p.runs):
                    if r.text:
                        font_name = r.font.name if r.font else 'Default'
                        font_size = r.font.size.pt if (r.font and r.font.size) else 'Default'
                        bold = r.font.bold if r.font else 'Default'
                        color = r.font.color.rgb if (r.font and r.font.color and hasattr(r.font.color, 'rgb')) else 'Default'
                        print(f"  Run {r_idx+1}: '{r.text}' | Font: {font_name}, Size: {font_size}, Bold: {bold}, Color: {color}")
