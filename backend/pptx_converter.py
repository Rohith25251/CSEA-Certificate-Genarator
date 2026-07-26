import io
import re
import base64
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

def get_color_hex(run_or_font):
    try:
        if run_or_font.font.color and run_or_font.font.color.type == 1:
            rgb = run_or_font.font.color.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
    except Exception:
        pass
    return "#000000"

def get_alignment_css(paragraph):
    align = paragraph.alignment
    if align == PP_ALIGN.CENTER:
        return "center"
    elif align == PP_ALIGN.RIGHT:
        return "right"
    elif align == PP_ALIGN.JUSTIFY:
        return "justify"
    return "left"

def convert_pptx_to_html_template(pptx_bytes: bytes) -> str:
    """Converts uploaded PPTX slide into responsive HTML template with exact text layout & placeholders."""
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        if not prs.slides:
            return ""

        slide = prs.slides[0]
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height

        # Slide width/height in px (Assuming 96 DPI / Emu conversion)
        width_px = int(slide_width_emu / 914400 * 96)
        height_px = int(slide_height_emu / 914400 * 96)

        shape_html_elements = []

        for shape in slide.shapes:
            left_pct = (shape.left / slide_width_emu) * 100 if slide_width_emu else 0
            top_pct = (shape.top / slide_height_emu) * 100 if slide_height_emu else 0
            width_pct = (shape.width / slide_width_emu) * 100 if slide_width_emu else 100
            height_pct = (shape.height / slide_height_emu) * 100 if slide_height_emu else 10

            # 1. Handle Picture shapes (Logos / Background images)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_bytes = image.blob
                    img_ext = image.ext
                    b64_str = base64.b64encode(img_bytes).decode('utf-8')
                    mime_type = f"image/{img_ext}"
                    img_html = f'''<div style="position: absolute; left: {left_pct:.2f}%; top: {top_pct:.2f}%; width: {width_pct:.2f}%; height: {height_pct:.2f}%; pointer-events: none;">
                        <img src="data:{mime_type};base64,{b64_str}" style="width: 100%; height: 100%; object-fit: contain;" alt="Embedded Element" />
                    </div>'''
                    shape_html_elements.append(img_html)
                    continue
                except Exception as img_err:
                    print(f"[PPTX Converter] Image extraction error: {img_err}")

            # 2. Handle Text Frame shapes
            if shape.has_text_frame:
                paragraphs_html = []
                for p in shape.text_frame.paragraphs:
                    align_css = get_alignment_css(p)
                    runs_html = []

                    for r in p.runs:
                        text = r.text
                        if not text:
                            continue

                        # Font properties
                        font_name = r.font.name or "Playfair Display, serif"
                        font_size_pt = r.font.size.pt if (r.font.size and hasattr(r.font.size, 'pt')) else 14
                        color_hex = get_color_hex(r)
                        is_bold = r.font.bold
                        is_italic = r.font.italic

                        style_parts = [
                            f"font-family: '{font_name}', sans-serif",
                            f"font-size: {font_size_pt:.1f}pt",
                            f"color: {color_hex}",
                            "display: inline-block"
                        ]

                        if is_bold:
                            style_parts.append("font-weight: bold")
                        if is_italic:
                            style_parts.append("font-style: italic")

                        run_span = f'<span style="{"; ".join(style_parts)}">{text}</span>'
                        runs_html.append(run_span)

                    if runs_html:
                        p_html = f'<p style="text-align: {align_css}; margin: 2px 0; line-height: 1.25;">{"".join(runs_html)}</p>'
                        paragraphs_html.append(p_html)

                if paragraphs_html:
                    container_html = f'''<div style="position: absolute; left: {left_pct:.2f}%; top: {top_pct:.2f}%; width: {width_pct:.2f}%; box-sizing: border-box; padding: 4px;">
                        {"".join(paragraphs_html)}
                    </div>'''
                    shape_html_elements.append(container_html)

        # Wrap in main container matching PowerPoint slide dimensions
        wrapper_html = f'''
<div style="position: relative; width: 1000px; height: 700px; max-width: 100%; background: #ffffff; border: 8px solid #082849; border-radius: 12px; overflow: hidden; box-sizing: border-box; margin: 0 auto; box-shadow: 0 20px 25px -5px rgba(0,0,0,0.1);">
    {"".join(shape_html_elements)}
</div>
        '''.strip()

        return wrapper_html
    except Exception as e:
        print(f"[PPTX Converter] Conversion error: {e}")
        return ""
