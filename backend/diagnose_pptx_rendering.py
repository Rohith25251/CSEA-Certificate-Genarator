import os
from pptx import Presentation
import win32com.client
import pythoncom

pptx_path = r"C:\Users\ROHITH P\Downloads\WORKSHOP.pptx"
prs = Presentation(pptx_path)

print("Analyzing PPTX TextFrames and word wrap properties...")
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            print(f"Shape: {shape.name} | WordWrap: {tf.word_wrap} | MarginLeft: {tf.margin_left} | MarginRight: {tf.margin_right}")
            for p in tf.paragraphs:
                print(f"  Paragraph: '{p.text}' | Font Size: {[r.font.size.pt if r.font and r.font.size else None for r in p.runs]}")
