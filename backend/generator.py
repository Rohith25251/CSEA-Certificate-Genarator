import os
import io
import re
import uuid
import base64
import tempfile
import datetime
import pandas as pd
from pptx import Presentation


# Attempt win32com and pythoncom imports for native PowerPoint COM rendering on Windows
has_win32com = False
pythoncom_lib = None
try:
    import win32com.client
    import pythoncom
    pythoncom_lib = pythoncom
    has_win32com = True
except Exception as err:
    if os.name == 'nt':
        print(f"[Generator] win32com/pythoncom import notice: {err}")
    has_win32com = False

# Attempt PyMuPDF for PDF page rasterization
has_fitz = False
try:
    import fitz
    has_fitz = True
except Exception:
    has_fitz = False

DEFAULT_CERTIFICATE_HTML = """
<div style="width: 1000px; height: 700px; padding: 45px; background: #ffffff; border: 12px solid #082849; border-radius: 12px; font-family: 'Playfair Display', serif; color: #1e293b; box-sizing: border-box; position: relative;">
  <div style="border: 2px solid #e2e8f0; height: 100%; padding: 40px; border-radius: 8px; text-align: center; display: flex; flex-direction: column; justify-content: space-between; box-sizing: border-box;">
    <div style="margin: 30px 0;">
      <h1 style="font-size: 32px; font-weight: 900; color: #082849;">CERTIFICATE OF PARTICIPATION</h1>
      <h2 style="font-size: 38px; font-weight: 800; color: #0284c7;"><<Name>></h2>
    </div>
  </div>
</div>
"""

is_vercel = os.environ.get('VERCEL') is not None or os.environ.get('AWS_LAMBDA_FUNCTION_NAME') is not None

if is_vercel:
    UPLOAD_DIR = os.path.join('/tmp', 'uploaded_templates')
else:
    UPLOAD_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), 'uploaded_templates'))
os.makedirs(UPLOAD_DIR, exist_ok=True)

SAVED_PPTX_PATH = os.path.join(UPLOAD_DIR, 'latest_template.pptx')
DOWNLOADS_WORKSHOP_PATH = r"C:\Users\ROHITH P\Downloads\WORKSHOP.pptx"

def get_active_pptx_template_path(event_id: str = None) -> str:
    """Returns path to active PPTX template, prioritizing downloading latest template from Supabase Storage bucket."""
    from supabase_client import download_latest_template_from_supabase
    
    if event_id:
        sp_path = download_latest_template_from_supabase(event_id)
        if sp_path and os.path.exists(sp_path):
            return sp_path

    sp_latest = download_latest_template_from_supabase()
    if sp_latest and os.path.exists(sp_latest):
        return sp_latest

    if event_id:
        event_local = os.path.join(UPLOAD_DIR, f"{event_id}_template.pptx")
        if os.path.exists(event_local):
            return event_local

    if os.path.exists(SAVED_PPTX_PATH):
        return SAVED_PPTX_PATH
    if os.path.exists(DOWNLOADS_WORKSHOP_PATH):
        return DOWNLOADS_WORKSHOP_PATH
    return ""

def save_uploaded_pptx_template(pptx_bytes: bytes) -> str:
    """Saves uploaded PPTX bytes to disk as latest_template.pptx."""
    with open(SAVED_PPTX_PATH, 'wb') as f:
        f.write(pptx_bytes)
    return SAVED_PPTX_PATH

def extract_placeholders_from_pptx(pptx_bytes: bytes):
    """Parses PPTX file, extracts text & placeholders, and saves template."""
    save_uploaded_pptx_template(pptx_bytes)
    
    prs = Presentation(io.BytesIO(pptx_bytes))
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    text_runs.append(p.text)
    
    combined_text = "\n".join(text_runs)
    placeholders = list(set(re.findall(r"<<\s*([^>]+?)\s*>>", combined_text)))
    
    return {
        "text": combined_text,
        "placeholders": placeholders,
        "template_saved": True
    }

from pptx.dml.color import RGBColor

def format_date_to_dd_mm_yyyy(val) -> str:
    if val is None:
        return ""
    if isinstance(val, (datetime.date, datetime.datetime)):
        return val.strftime("%d\u2011%m\u2011%Y")
    
    val_str = str(val).strip()
    
    # Check if format is yyyy-mm-dd (with optional T... or time)
    match1 = re.match(r'^(\d{4})[-/](\d{2})[-/](\d{2})(?:\s+.*|T.*)?$', val_str)
    if match1:
        y, m, d = match1.groups()
        return f"{d}\u2011{m}\u2011{y}"
        
    # Check if format is already dd-mm-yyyy or dd/mm/yyyy
    match2 = re.match(r'^(\d{2})[-/](\d{2})[-/](\d{4})$', val_str)
    if match2:
        d, m, y = match2.groups()
        return f"{d}\u2011{m}\u2011{y}"
        
    return val_str

def build_dynamic_replacements(row: dict, extra: dict = None) -> dict:
    """
    Builds a universal dynamic replacement dictionary from ANY arbitrary Excel column header,
    handling trailing/leading spaces, casing, underscore/space variations, and common aliases.
    """
    rep = {}
    if extra:
        rep.update(extra)

    # 1. First pass: Add all original Excel column key-value pairs and clean variations
    for raw_k, v in row.items():
        if raw_k is None:
            continue
        val_str = str(v).strip() if v is not None else ""
        clean_k = str(raw_k).strip()

        rep[raw_k] = val_str
        rep[clean_k] = val_str
        rep[clean_k.lower()] = val_str
        rep[clean_k.upper()] = val_str
        rep[clean_k.title()] = val_str
        rep[clean_k.replace('_', ' ')] = val_str
        rep[clean_k.replace(' ', '_')] = val_str
        rep[clean_k.replace('_', ' ').lower()] = val_str

    # 2. Extract standard fields (Name, Roll No, Email, Date, College, Event) if present
    std_name = None
    std_roll = None
    std_email = None

    for k, v in row.items():
        if not k or v is None:
            continue
        ck = str(k).strip().lower().replace('_', ' ')
        cv = str(v).strip()

        if ('name' in ck or 'student' in ck or 'participant' in ck) and not std_name and cv:
            std_name = cv
        if ('roll' in ck or 'register' in ck or 'reg' in ck or 'id' in ck or 'number' in ck) and not std_roll and cv:
            std_roll = cv
        if ('mail' in ck or 'email' in ck) and '@' in cv:
            std_email = cv

    # 3. Register standard aliases for PowerPoint templates with varied token naming
    if std_name:
        for alias in ["Name", "Name ", "name", "Full Name", "Full Name ", "Student Name", "Participant Name", "PARTICIPANT_NAME"]:
            if alias not in rep:
                rep[alias] = std_name

    if std_roll:
        for alias in [
            "RollNumber", "Roll Number", "Roll Number ", "Roll No", "Roll No ", "Roll_Number", "Roll_No", 
            "RollNo", "Register No", "Register No ", "RegisterNo", "Register_No", "registration_no", 
            "Registration Number", "REGISTRATION_NO", "Reg No", "Reg_No", "Student ID", "StudentID"
        ]:
            rep[alias] = std_roll
            rep[alias.lower()] = std_roll
            rep[alias.upper()] = std_roll
            rep[alias.title()] = std_roll

    if std_email:
        for alias in ["Email", "Mail id", "Mail ID", "College mail id", "student_email"]:
            if alias not in rep:
                rep[alias] = std_email

    # 4. Format any date-like values in the replacements dictionary to dd-mm-yyyy with non-breaking hyphens
    for k, v in list(rep.items()):
        if isinstance(k, str) and ('date' in k.lower() or k.lower() == 'date'):
            rep[k] = format_date_to_dd_mm_yyyy(v)

    return rep


def replace_tokens_in_pptx_slide(slide, replacements: dict):
    """Replaces <<Placeholder>> tokens inside PowerPoint slides while preserving 100% of font family, font size, bold, italic, and colors."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        try:
            shape.text_frame.word_wrap = True
        except Exception:
            pass

        # Normalize fonts and disable underlines slide-wide
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font:
                    r.font.underline = False
                    if r.font.name:
                        name_lower = r.font.name.lower()
                        if "times" in name_lower:
                            r.font.name = "Times New Roman"
                            if "bold" in name_lower:
                                r.font.bold = True
                        elif "playfair" in name_lower:
                            r.font.name = "Playfair Display"
                            if "bold" in name_lower:
                                r.font.bold = True

        for p in shape.text_frame.paragraphs:
            full_text = p.text
            if not full_text or '<<' not in full_text:
                continue

            # Find all <<token>> matches with character start/end indices in full_text
            matches = []
            for m in re.finditer(r"<<\s*([^>]+?)\s*>>", full_text):
                token_clean = m.group(1).strip()
                val = None
                for k in [m.group(1), token_clean, token_clean.lower(), token_clean.upper(), token_clean.title(), token_clean.replace('_', ' '), token_clean.replace(' ', '_')]:
                    if k in replacements:
                        val = replacements[k]
                        break
                if val is not None:
                    val_str = str(val)
                    # If the placeholder is followed immediately by an alphanumeric character (like "o" in "organized"),
                    # append a space to prevent breaking words and line wraps.
                    if m.end() < len(full_text) and full_text[m.end()].isalnum():
                        val_str += " "
                    matches.append((m.start(), m.end(), val_str))

            if not matches:
                continue

            # Map each character index in full_text to (run_index, char_index_within_run)
            char_map = []
            for r_idx, r in enumerate(p.runs):
                for c_idx in range(len(r.text)):
                    char_map.append((r_idx, c_idx))

            # Perform replacement from right to left so indices remain valid
            for start, end, val_str in reversed(matches):
                if start >= len(char_map) or end > len(char_map):
                    continue

                start_run_idx, start_char_offset = char_map[start]
                end_run_idx, end_char_offset = char_map[end - 1]

                if start_run_idx == end_run_idx:
                    # Token is inside a single run
                    r = p.runs[start_run_idx]
                    r.text = r.text[:start_char_offset] + val_str + r.text[end_char_offset + 1:]
                else:
                    # Token spans across multiple runs
                    r_first = p.runs[start_run_idx]
                    r_first.text = r_first.text[:start_char_offset] + val_str

                    for mid_idx in range(start_run_idx + 1, end_run_idx):
                        p.runs[mid_idx].text = ""

                    r_last = p.runs[end_run_idx]
                    r_last.text = r_last.text[end_char_offset + 1:]

def replace_html_tokens(html_content: str, replacements: dict) -> str:
    """Replaces <<placeholder>> and {{placeholder}} in HTML content with replacement values (case-insensitive)."""
    rendered = html_content
    for k, v in replacements.items():
        val = str(v) if v is not None else ""
        escaped_k = re.escape(str(k).strip())
        pattern = r"(<<\s*" + escaped_k + r"\s*>>|\{\{\s*" + escaped_k + r"\s*\}\})"
        rendered = re.sub(pattern, val, rendered, flags=re.IGNORECASE)
    return rendered

def generate_single_native_pdf(pptx_template_path: str, replacements: dict, output_pdf_path: str) -> bool:
    """Modifies PPTX template with student replacements and converts directly to PDF using PowerPoint COM with CoInitialize."""
    target_pptx = pptx_template_path if (pptx_template_path and os.path.exists(pptx_template_path)) else get_active_pptx_template_path()
    
    if not target_pptx or not os.path.exists(target_pptx):
        print(f"[Generator] PPTX template not found at {target_pptx}")
        return False

    temp_dir = tempfile.mkdtemp()
    temp_pptx = os.path.join(temp_dir, f"temp_{uuid.uuid4().hex[:6]}.pptx")

    co_initialized = False
    try:
        prs = Presentation(target_pptx)
        slide = prs.slides[0]
        replace_tokens_in_pptx_slide(slide, replacements)
        prs.save(temp_pptx)

        if has_win32com:
            if pythoncom_lib:
                pythoncom_lib.CoInitialize()
                co_initialized = True

            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            pres = ppt_app.Presentations.Open(os.path.abspath(temp_pptx), WithWindow=False)
            pres.SaveAs(os.path.abspath(output_pdf_path), 32) # 32 = ppSaveAsPDF
            pres.Close()
            ppt_app.Quit()

            if co_initialized:
                pythoncom_lib.CoUninitialize()

            print(f"[Generator] Successfully rendered native PowerPoint PDF at {output_pdf_path}")
            return True
        else:
            print("[Generator] win32com unavailable. Falling back to high-fidelity ReportLab PDF rendering.")
            try:
                import io
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                from pptx.enum.text import PP_ALIGN
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import A4, landscape
                from reportlab.platypus import Paragraph, Frame
                from reportlab.lib.styles import ParagraphStyle
                from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
                from reportlab.lib.utils import ImageReader
                from reportlab.lib.colors import HexColor
                from svglib.svglib import svg2rlg

                # 1. Clean replacements to remove any Unicode non-breaking hyphens
                cleaned_replacements = {}
                for k, v in replacements.items():
                    if isinstance(v, str):
                        cleaned_replacements[k] = v.replace('\u2011', '-')
                    else:
                        cleaned_replacements[k] = v

                # Use the modified temporary PPTX
                prs_temp = Presentation(temp_pptx)
                slide_temp = prs_temp.slides[0]

                # A4 Page dimensions in points
                A4_w, A4_h = landscape(A4) # 841.89 x 595.27

                os.makedirs(os.path.dirname(os.path.abspath(output_pdf_path)), exist_ok=True)
                pdf_canvas = canvas.Canvas(output_pdf_path, pagesize=(A4_w, A4_h))

                slide_w = prs_temp.slide_width
                slide_h = prs_temp.slide_height
                scale_x = A4_w / slide_w
                scale_y = A4_h / slide_h

                for shape_idx, shape in enumerate(slide_temp.shapes):
                    x = shape.left * scale_x
                    y = A4_h - (shape.top + shape.height) * scale_y
                    w = shape.width * scale_x
                    h = shape.height * scale_y

                    # Picture shape OR shape with picture fill
                    has_picture_fill = False
                    img_bytes = None
                    is_svg = False
                    if hasattr(shape, 'fill') and shape.fill and shape.fill.type == 6: # MSO_FILL.PICTURE
                        try:
                            # Use universal xpath to match r:embed on standard blips or svgBlips
                            rIds = shape.fill._xPr.xpath('.//@r:embed')
                            rId = rIds[0] if rIds else None
                            if rId:
                                part = slide_temp._part.related_part(rId)
                                img_bytes = part.blob
                                has_picture_fill = True
                                if hasattr(part, 'partname') and str(part.partname).lower().endswith('.svg'):
                                    is_svg = True
                                elif img_bytes.startswith(b'<svg') or img_bytes.startswith(b'<?xml') or b'<svg' in img_bytes[:200]:
                                    is_svg = True
                        except Exception:
                            pass

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or has_picture_fill:
                        try:
                            if not img_bytes and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                img_bytes = shape.image.blob
                                if hasattr(shape.image, 'filename') and str(shape.image.filename).lower().endswith('.svg'):
                                    is_svg = True
                            
                            if img_bytes:
                                if is_svg:
                                    drawing = svg2rlg(io.BytesIO(img_bytes))
                                    sx = w / drawing.width
                                    sy = h / drawing.height
                                    drawing.scale(sx, sy)
                                    drawing.drawOn(pdf_canvas, x, y)
                                else:
                                    img_io = io.BytesIO(img_bytes)
                                    img_reader = ImageReader(img_io)
                                    pdf_canvas.drawImage(img_reader, x, y, w, h, mask='auto')
                        except Exception as img_err:
                            print(f"[Generator] Fallback error rendering shape [{shape_idx}]: {img_err}")

                    # AutoShape line divider (height is 0)
                    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.height == 0:
                        try:
                            if shape.line and shape.line.color and hasattr(shape.line.color, 'rgb') and shape.line.color.rgb:
                                rgb = shape.line.color.rgb
                                color_hex = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                lw = (shape.line.width / 12700.0) if (shape.line.width) else 1.0
                                
                                pdf_canvas.setStrokeColor(HexColor(color_hex))
                                pdf_canvas.setLineWidth(lw)
                                pdf_canvas.line(x, y + h, x + w, y + h)
                        except Exception:
                            pass

                    # Text Box
                    if shape.has_text_frame:
                        try:
                            story = []
                            has_text = False

                            for paragraph in shape.text_frame.paragraphs:
                                p_text = ""
                                max_font_size = 12

                                for run in paragraph.runs:
                                    text = run.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                    text = text.replace('\u2011', '-')
                                    if not text.strip():
                                        p_text += text
                                        continue

                                    has_text = True
                                    style_start = ""
                                    style_end = ""

                                    if run.font.bold:
                                        style_start += "<b>"
                                        style_end = "</b>" + style_end
                                    if run.font.italic:
                                        style_start += "<i>"
                                        style_end = "</i>" + style_end

                                    font_name = run.font.name or "Helvetica"
                                    if "times" in font_name.lower() or "playfair" in font_name.lower():
                                        if run.font.bold and run.font.italic:
                                            rl_font = "Times-BoldItalic"
                                        elif run.font.bold:
                                            rl_font = "Times-Bold"
                                        elif run.font.italic:
                                            rl_font = "Times-Italic"
                                        else:
                                            rl_font = "Times-Roman"
                                    else:
                                        if run.font.bold and run.font.italic:
                                            rl_font = "Helvetica-BoldOblique"
                                        elif run.font.bold:
                                            rl_font = "Helvetica-Bold"
                                        elif run.font.italic:
                                            rl_font = "Helvetica-Oblique"
                                        else:
                                            rl_font = "Helvetica"

                                    size_pt = run.font.size.pt if (run.font.size and hasattr(run.font.size, 'pt')) else 14
                                    if size_pt > max_font_size:
                                        max_font_size = size_pt

                                    color_hex = "#000000"
                                    try:
                                        if run.font.color and run.font.color.type == 1:
                                            rgb = run.font.color.rgb
                                            color_hex = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                    except Exception:
                                        pass

                                    style_start += f'<font name="{rl_font}" size="{size_pt:.1f}" color="{color_hex}">'
                                    style_end = "</font>" + style_end

                                    p_text += f"{style_start}{text}{style_end}"

                                if has_text:
                                    align = TA_LEFT
                                    if paragraph.alignment == PP_ALIGN.CENTER:
                                        align = TA_CENTER
                                    elif paragraph.alignment == PP_ALIGN.RIGHT:
                                        align = TA_RIGHT
                                    elif paragraph.alignment == PP_ALIGN.JUSTIFY:
                                        align = TA_JUSTIFY

                                    leading = max_font_size * 1.25
                                    p_style = ParagraphStyle(
                                        name=f"style_{uuid.uuid4().hex[:6]}",
                                        alignment=align,
                                        leading=leading
                                    )
                                    story.append(Paragraph(p_text, p_style))

                            if story:
                                f = Frame(x, y, w, h, leftPadding=2, rightPadding=2, topPadding=2, bottomPadding=2, id=None)
                                f.addFromList(story, pdf_canvas)
                        except Exception as txt_err:
                            print(f"[Generator] Fallback error rendering text shape [{shape_idx}]: {txt_err}")

                pdf_canvas.showPage()
                pdf_canvas.save()
                print(f"[Generator] Successfully rendered high-fidelity fallback PDF using ReportLab at {output_pdf_path}")
                return True
            except Exception as fallback_err:
                print(f"[Generator] Fallback ReportLab rendering failed: {fallback_err}")
                return False
    except Exception as e:
        print(f"[Generator] Native PPTX to PDF error: {e}")
        if co_initialized and pythoncom_lib:
            try:
                pythoncom_lib.CoUninitialize()
            except Exception:
                pass
        return False

def pdf_to_base64_png(pdf_path: str) -> str:
    """Renders page 1 of PDF file to Base64 PNG data URL."""
    if not has_fitz or not os.path.exists(pdf_path):
        return ""
    try:
        doc = fitz.open(pdf_path)
        page = doc[0]
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("png")
        b64 = base64.b64encode(img_bytes).decode('utf-8')
        return f"data:image/png;base64,{b64}"
    except Exception as e:
        print(f"[Generator] PDF to PNG error: {e}")
        return ""

def parse_excel_dataframe(file_bytes: bytes, filename: str):
    """Parses Excel bytes into records dictionary."""
    if filename.endswith('.csv'):
        df = pd.read_csv(io.BytesIO(file_bytes))
    else:
        df = pd.read_excel(io.BytesIO(file_bytes))

    df = df.fillna('')
    headers = list(df.columns)
    rows = df.to_dict(orient='records')
    return {
        "headers": headers,
        "rows": rows,
        "total_rows": len(rows)
    }
