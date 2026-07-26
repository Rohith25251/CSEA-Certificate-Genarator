import os
import re
from pptx import Presentation
from pptx.dml.color import RGBColor

def build_dynamic_replacements(row: dict, extra: dict = None) -> dict:
    """
    Builds a flexible replacement map from any arbitrary Excel row,
    populating normalized key variations (lowercase, stripped, space/underscore agnostic).
    """
    rep = {}
    if extra:
        rep.update(extra)
        
    for k, v in row.items():
        if k is None:
            continue
        val_str = str(v) if v is not None else ""
        raw_key = str(k)
        clean_key = raw_key.strip()
        
        # Add original key and cleaned variations
        rep[raw_key] = val_str
        rep[clean_key] = val_str
        rep[clean_key.lower()] = val_str
        rep[clean_key.upper()] = val_str
        rep[clean_key.title()] = val_str
        
        # Add space <-> underscore variations
        rep[clean_key.replace('_', ' ')] = val_str
        rep[clean_key.replace(' ', '_')] = val_str

    return rep

def replace_tokens_in_pptx_slide_universal(slide, replacements: dict):
    """
    Universally replaces ALL <<placeholder>> tokens found in any slide shape,
    matching arbitrary Excel column names dynamically.
    """
    blue_color = RGBColor(0x1A, 0x39, 0x87)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        try:
            shape.text_frame.word_wrap = True
        except Exception:
            pass

        for p in shape.text_frame.paragraphs:
            if not p.text:
                continue

            for r in p.runs:
                if not r.text:
                    continue

                # Enforce Bold & Blue (#1A3987) for DATASET TO DECISION or event_name
                if 'DATASET TO DECISION' in r.text or '<<event_name>>' in r.text.lower() or '<<event name>>' in r.text.lower():
                    if r.font:
                        r.font.bold = True
                        r.font.color.rgb = blue_color

                if '<<' in r.text:
                    run_text = r.text
                    # Find all placeholders in this run like <<Token>>
                    found_tokens = re.findall(r"<<\s*([^>]+?)\s*>>", run_text)
                    
                    for token in found_tokens:
                        token_clean = token.strip()
                        # Search for token match in dynamic replacements map
                        matched_val = None
                        for key_var in [token_clean, token_clean.lower(), token_clean.upper(), token_clean.title(), token_clean.replace('_', ' '), token_clean.replace(' ', '_')]:
                            if key_var in replacements:
                                matched_val = replacements[key_var]
                                break
                        
                        if matched_val is not None:
                            pattern = re.compile(rf"<<\s*{re.escape(token)}\s*>>", re.IGNORECASE)
                            run_text = pattern.sub(str(matched_val), run_text)

                    r.text = run_text

                # Post-check for event name bold/blue styling
                event_val = str(replacements.get('event_name') or replacements.get('Event Name') or 'DATASET TO DECISION').strip()
                if event_val and event_val.lower() in r.text.lower() and len(event_val) > 2:
                    if r.font:
                        r.font.bold = True
                        r.font.color.rgb = blue_color

# Test execution with custom arbitrary Excel columns
custom_excel_row = {
    "Participant_Name ": "Rohith P",
    "Registration_No": "25CSR999",
    "Specialization": "Artificial Intelligence",
    "Grade": "Distinction",
    "Mail ID": "rohith@kongu.edu",
    "Event_Title": "ADVANCED AGENTIC CODING WORKSHOP",
    "Date ": "25.07.2026"
}

rep_map = build_dynamic_replacements(custom_excel_row, {"event_name": "ADVANCED AGENTIC CODING WORKSHOP"})
print("Generated Dynamic Replacement Keys:", list(rep_map.keys())[:15])

pptx_path = os.path.abspath(os.path.join(os.path.dirname(__file__), 'uploaded_templates', 'latest_template.pptx'))
prs = Presentation(pptx_path)
slide = prs.slides[0]

replace_tokens_in_pptx_slide_universal(slide, rep_map)

print("\n--- TEST REPLACEMENT RESULTS ---")
for shape in slide.shapes:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            if p.text.strip():
                print(f"Shape Paragraph: '{p.text}'")
